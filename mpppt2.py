import os
import time

import comtypes.client
import cv2
import imagehash
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


def extract_and_crop_images(pptx_path, output_dir):
    image_paths = []
    # 절대 경로를 사용하여 경로 문제를 방지합니다.
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)
    
    # 출력 폴더가 존재하지 않으면 생성합니다.
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    
    # 프레젠테이션 열기
    presentation = powerpoint.Presentations.Open(pptx_path)
    
    # 슬라이드마다 개별 이미지로 저장
    for i, slide in enumerate(presentation.Slides):
        # 파일 이름 지정
        slide_name = f"slide_{i + 1}.png"
        slide_path = os.path.join(output_dir, slide_name)
        image_paths.append(slide_path)
        # 슬라이드 내보내기
        slide.Export(slide_path, "PNG")
    
    # 프레젠테이션 닫기 및 PowerPoint 종료
    presentation.Close()
    powerpoint.Quit()
    
    # 이미지 크롭 후 원본 이미지 삭제
    cropped_image_paths = []
    for image_path in image_paths:
        image = Image.open(image_path)
        cropped_image = image.crop((0, 30, 1080, 600))
        cropped_image_path = image_path.replace(".png", "_cropped.png")
        cropped_image.save(cropped_image_path)
        cropped_image_paths.append(cropped_image_path)
        os.remove(image_path)
                
    return cropped_image_paths

def calculate_similarity(img1, img2):
    hash1 = imagehash.average_hash(Image.open(img1))
    hash2 = imagehash.average_hash(Image.open(img2))
    return hash1 - hash2

def find_best_frame(video_path, image_path):
    try:
        video = cv2.VideoCapture(video_path)
        frame_rate = video.get(cv2.CAP_PROP_FPS)
        success, frame = video.read()
        frame_id = 0
        min_similarity = float('inf')
        best_frame_path = None

        while success:
            if frame_id % int(frame_rate) == 0:
                frame_path = f"frame_{frame_id}.jpg"
                cv2.imwrite(frame_path, frame)
                
                similarity = calculate_similarity(image_path, frame_path)
                if similarity < min_similarity:
                    min_similarity = similarity
                    if best_frame_path:
                        os.remove(best_frame_path)
                    best_frame_path = frame_path
                else:
                    time.sleep(0.1)  # 지연 시간 추가
                    os.remove(frame_path)
                
            success, frame = video.read()
            frame_id += 1

        video.release()
    except Exception as e:
        pass
    return best_frame_path, min_similarity

def add_image_to_ppt(ppt, img_path, similarity, slide_index):
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    
    # 제목을 수동으로 추가
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_text_frame = title_shape.text_frame
        title_text_frame.text = f"슬라이드 {slide_index + 1}의 유사도 결과"
    else:
        title_shape.text = f"슬라이드 {slide_index + 1}의 유사도 결과"

    slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(3))

    text_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(5), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = f"유사도: {similarity}"

def main(ppt_path, video_path, output_dir):
    # Step 1: PPT 슬라이드에서 이미지 추출 및 크롭
    image_paths = extract_and_crop_images(ppt_path, output_dir)
    
    # 새로운 프레젠테이션 객체 생성
    result_ppt = Presentation()
    
    # Step 2: 동영상에서 각 슬라이드에 대한 유사한 장면 찾기 및 결과 PPT에 추가
    for i, image_path in enumerate(image_paths):
        best_frame, min_similarity = find_best_frame(video_path, image_path)
        
        # Step 3: 캡쳐본을 결과 PPT에 추가 및 유사도 표시
        add_image_to_ppt(result_ppt, best_frame, min_similarity, i)
        
        # 캡쳐본 파일 이름 변경
        if best_frame:
            os.rename(best_frame, f"frame_{i}.jpg")
    
    # 결과 PPT 저장
    result_ppt_path = "result_presentation.pptx"
    result_ppt.save(result_ppt_path)
    
    print(f"결과 파일: {result_ppt_path}")

if __name__ == "__main__":
    ppt_path = "C://Users//wpghk//Desktop//automatic//test.pptx"
    video_path = "C://Users//wpghk//Desktop//automatic//(종편본)관광창업론_01주차_02_학습내용_수정07.mp4"
    image_directory = "C://Users//wpghk//Desktop//automatic//slides"
    
    main(ppt_path, video_path, image_directory)
