import os
import pickle
import re
import shutil
from concurrent.futures import ThreadPoolExecutor

import comtypes.client
import cv2
import easyocr
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# EasyOCR 설정
reader = easyocr.Reader(['en', 'ko'], gpu=True)  # 영어와 한국어 설정

def pptx_to_images(pptx_path, output_dir):
    # 절대 경로를 사용하여 경로 문제를 방지합니다.
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)
    
    # 출력 폴더가 존재하지 않으면 생성합니다.
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    
    # 프레젠테이션 열기
    presentation = powerpoint.Presentations.Open(pptx_path)
    total_slides = len(presentation.slides)
    
    # 슬라이드마다 개별 이미지로 저장
    for i, slide in enumerate(presentation.Slides):
        slide_num = str(i + 1).zfill(len(str(total_slides)))
        # 파일 이름 지정
        slide_name = f"slide_{slide_num}.png"
        slide_path = os.path.join(output_dir, slide_name)
        
        # 슬라이드 내보내기
        slide.Export(slide_path, "PNG")
    
    # 프레젠테이션 닫기 및 PowerPoint 종료
    presentation.Close()
    powerpoint.Quit()
    
def extract_text_from_image(image):
    result = reader.readtext(np.array(image))
    text = ' '.join([item[1] for item in result])
    return text

def extract_text_from_ppt(directory, ocr_area=(0, 90, 1080, 600)):
    slide_texts = []
    
    image_files = sorted(
        [f for f in os.listdir(directory) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff'))]
    )

    for image_file in image_files:
        image_path = os.path.join(directory, image_file)
        image = Image.open(image_path)
        cropped_image = image.crop(ocr_area)
        text = extract_text_from_image(cropped_image)
        # 예) # 뒤에 숫자가 오는 글자는 빈 글자로 replace
        text = re.sub(r'\#\d+', '', text)
        
        slide_texts.append(text)
    
    # slide_texts text 파일로 저장
    with open('slide_texts.txt', 'w') as f:
        for text in slide_texts:
            f.write(text + '\n')
    
    return slide_texts

def process_frame(frame, scale_factor=0.5):
    resized_frame = cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor)
    cropped_frame = resized_frame[60:600, 0:1080]
    gray = cv2.cvtColor(cropped_frame, cv2.COLOR_BGR2GRAY)
    result = reader.readtext(gray)
    text = ' '.join([item[1] for item in result])
    return text

def cache_frame_texts(video_path, frame_skip, scale_factor=0.5, cache_file='frame_texts_cache.pkl'):
    cap = cv2.VideoCapture(video_path)
    frame_texts = []
    frame_images = []
    frame_indices = []
    frame_count = 30
    
    while True:
        success, frame = cap.read()
        if not success:
            break
        if frame_count % frame_skip == 0:
            text = process_frame(frame, scale_factor)
            resized_frame = cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor)
            frame_texts.append(text)
            frame_images.append(resized_frame)
            frame_indices.append(frame_count)
        frame_count += 1
    
    # frame_texts 파일로 저장
    with open('frame_texts.txt', 'w') as f:
        for text in frame_texts:
            f.write(text + '\n')
    
    with open(cache_file, 'wb') as f:
        pickle.dump((frame_texts, frame_images, frame_indices), f)

def load_cached_frame_texts(cache_file='frame_texts_cache.pkl'):
    with open(cache_file, 'rb') as f:
        return pickle.load(f)
    
def process_video_and_annotate_ppt(ppt_path, video_path, image_directory, frame_skip=60, cache_file='frame_texts_cache.pkl'):
    slide_texts = extract_text_from_ppt(image_directory)
    
    if not os.path.exists(cache_file):
        cache_frame_texts(video_path, frame_skip, cache_file=cache_file)
    
    frame_texts, frame_images, frame_indices = load_cached_frame_texts(cache_file)
    
    vectorizer = TfidfVectorizer().fit_transform(slide_texts + frame_texts)
    vectors = vectorizer.toarray()
    slide_vectors = vectors[:len(slide_texts)]
    frame_vectors = vectors[len(slide_texts):]
    similarities = cosine_similarity(slide_vectors, frame_vectors)
    
    most_similar_frames = similarities.argmax(axis=1)
    captured_images = [frame_images[i] for i in most_similar_frames]
    
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        max_sim_val = similarities[i][most_similar_frames[i]]
        img_path = os.path.join(image_directory, f"most_similar_frame_{i + 1}.png")
        cv2.imwrite(img_path, captured_images[i])
        left = Inches(0.5)
        top = Inches(5)
        # 가장 비슷한 사진 배치
        pic = slide.shapes.add_picture(img_path, Inches(-9.6), Inches(0.3), height=Inches(5), width=Inches(9.5))
        
        # 유사도를 표시하는 텍스트 박스 추가
        text_box = slide.shapes.add_textbox(Inches(11.3), top - Inches(0.5), Inches(2), Inches(0.5))
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.text = f"유사도 : {max_sim_val:.2f}"
        p.font.size = Pt(16)
        # 텍스트 가운데 정렬
        p.alignment = PP_ALIGN.CENTER

        # 텍스트 상자의 배경색을 노란색으로 설정
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 0)  # 노란색 (RGB 값)
        
        # 유사도가 0.4 이하일 경우 슬라이드 배경을 빨간색으로 설정
        if max_sim_val <= 0.25:
            # 레드카드 생성 왼쪽부터 왼쪽에서 멀어질 크기, 위에서 멀어질 크기, 너비, 높이
            rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(2.3), Inches(2), Inches(2.1))
            rectangle.fill.solid()
            rectangle.fill.fore_color.rgb = RGBColor(255, 0, 0)  # 빨간색 (RGB 값)
        
    prs.save("annotated_presentation.pptx")

# 실행 예제
ppt_path = "C://Users//wpghk//Desktop//automatic_text//test2.pptx"
image_directory = "C://Users//wpghk//Desktop//automatic_text//slides"
video_path = "C://Users//wpghk//Desktop//automatic_text//(종편본)관광창업론_01주차_01_학습목표_수정07.mp4"

#ppt_path = "C://Users//wpghk//Desktop//automatic_text//test.pptx"
#image_directory = "C://Users//wpghk//Desktop//automatic_text//slides"
#video_path = "C://Users//wpghk//Desktop//automatic_text//(종편본)관광창업론_01주차_02_학습내용_수정07.mp4"

pptx_to_images(ppt_path, image_directory)
process_video_and_annotate_ppt(ppt_path, video_path, image_directory)

# 캐시 파일 삭제
os.remove('frame_texts_cache.pkl')

# 캐시 파일이 삭제되었는지 확인
os.path.exists('frame_texts_cache.pkl')

# image_directory 삭제
shutil.rmtree(image_directory)
