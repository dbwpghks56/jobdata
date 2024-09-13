import os
import pickle
import re
import shutil
import urllib.parse
import uuid
from pathlib import Path

import comtypes.client
import cv2
import easyocr
import numpy as np
import pythoncom
from flask import Flask, jsonify, request, send_file
from flask_cors import CORS
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

app = Flask(__name__)
CORS(app)

# EasyOCR 설정
reader = easyocr.Reader(['en', 'ko'], gpu=True)  # 영어와 한국어 설정

def pptx_to_images(pptx_path, output_dir):
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    pythoncom.CoInitialize()
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(pptx_path)
    total_slides = len(presentation.slides)
    
    for i, slide in enumerate(presentation.Slides):
        slide_num = str(i + 1).zfill(len(str(total_slides)))
        slide_name = f"slide_{slide_num}.png"
        slide_path = os.path.join(output_dir, slide_name)
        slide.Export(slide_path, "PNG")
    
    presentation.Close()
    powerpoint.Quit()
    pythoncom.CoUninitialize()

def extract_text_from_image(image):
    result = reader.readtext(np.array(image))
    text = ' '.join([item[1] for item in result])
    return text

def extract_text_from_ppt(directory, ocr_area=(0, 90, 1080, 600)):
    """
    Extract text from PPT slide images within the specified OCR area.
    """
    slide_texts = []
    
    image_files = sorted(
        [f for f in os.listdir(directory) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff'))]
    )

    for image_file in image_files:
        image_path = os.path.join(directory, image_file)
        image = Image.open(image_path)
        cropped_image = image.crop(ocr_area)  # Use the OCR area specified
        cropped_image.save(image_path)  # Save the cropped image
        text = extract_text_from_image(cropped_image)
        text = re.sub(r'\#\d+', '', text)
        
        slide_texts.append(text)
        
    return slide_texts

def process_frame(frame, scale_factor=0.3):
    resized_frame = cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor)
    cropped_frame = resized_frame[60:600, 0:1080]
    gray = cv2.cvtColor(cropped_frame, cv2.COLOR_BGR2GRAY)
    result = reader.readtext(gray)
    text = ' '.join([item[1] for item in result])
    return text

def cache_frame_texts(video_path, frame_skip, scale_factor=0.3, cache_file='frame_texts_cache.pkl'):
    cap = cv2.VideoCapture(video_path)
    frame_texts = []
    frame_images = []
    frame_indices = []
    frame_count = 0
    
    while True:
        success, frame = cap.read()
        if not success:
            break
        if frame_count % frame_skip == 0:
            text = process_frame(frame, scale_factor)
            resized_frame = cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor)
            frame_texts.append(text)
            frame_images.append(resized_frame)
            frame_indices.append(frame_count)
        frame_count += 1
    
    with open(cache_file, 'wb') as f:
        pickle.dump((frame_texts, frame_images, frame_indices), f)

def load_cached_frame_texts(cache_file='frame_texts_cache.pkl'):
    with open(cache_file, 'rb') as f:
        return pickle.load(f)

def process_video_and_annotate_ppt(ppt_path, video_path, image_directory, output_path, ocr_area, frame_skip=30, cache_file='frame_texts_cache.pkl'):
    slide_texts = extract_text_from_ppt(image_directory, ocr_area=ocr_area)
    
    if not os.path.exists(cache_file):
        cache_frame_texts(video_path, frame_skip, cache_file=cache_file)
    
    frame_texts, frame_images, frame_indices = load_cached_frame_texts(cache_file)
    
    vectorizer = TfidfVectorizer().fit_transform(slide_texts + frame_texts)
    vectors = vectorizer.toarray()
    slide_vectors = vectors[:len(slide_texts)]
    frame_vectors = vectors[len(slide_texts):]
    similarities = cosine_similarity(slide_vectors, frame_vectors)
    
    most_similar_frames = similarities.argmax(axis=1)
    captured_images = [frame_images[i] for i in most_similar_frames]
    
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        max_sim_val = similarities[i][most_similar_frames[i]]
        img_path = image_directory + f"//most_similar_frame_{i + 1}.png"
        cv2.imwrite(img_path, captured_images[i])
        left = Inches(0.5)
        top = Inches(5)
        pic = slide.shapes.add_picture(img_path, Inches(-9.6), Inches(0.3), height=Inches(5), width=Inches(9.5))
        
        text_box = slide.shapes.add_textbox(Inches(11.3), top - Inches(0.5), Inches(2), Inches(0.5))
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.text = f"유사도 : {max_sim_val:.2f}"
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 0)
        
        if max_sim_val <= 0.25:
            rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(2.3), Inches(2), Inches(2.1))
            rectangle.fill.solid()
            rectangle.fill.fore_color.rgb = RGBColor(255, 0, 0)
        
    prs.save(output_path)
    return output_path

@app.route('/process', methods=['POST'])
def process():
    ppt_file = request.files['ppt_file']
    video_file = request.form.get('video_file_path')
    left = int(request.form.get('left', 0))
    top = int(request.form.get('top', 90))
    right = int(request.form.get('right', 1080))
    bottom = int(request.form.get('bottom', 600))
    
    ppt_path = os.path.join(os.getcwd(), ppt_file.filename)
    video_path = video_file
    
    ppt_file.save(ppt_path)
    
    unique = str(uuid.uuid4()).split('-')[0]
    image_directory = os.path.join(os.getcwd(), "slides_" + unique)
    output_path = os.path.join(os.getcwd(), "검수_" + os.path.basename(ppt_path))
    
    try:
        pptx_to_images(ppt_path, image_directory)
        output_file = process_video_and_annotate_ppt(ppt_path, video_path, image_directory, output_path, ocr_area=(left, top, right, bottom), cache_file=f'frame_texts_cache_{unique}.pkl')
        
        # Cleanup
        os.remove(f'frame_texts_cache_{unique}.pkl')
        shutil.rmtree(image_directory) 
        os.remove(ppt_path)
        
        return send_file(output_file, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/video_list', methods=['GET'])
def list_video_files_in_directory():
    video_extensions = ['.mp4', '.avi', '.mov', '.mkv', '.wmv']
    path = Path('uploads')
    fileList = []
    fileData = {}
    if path.exists() and path.is_dir():
        video_files = [file for file in path.iterdir() if file.suffix.lower() in video_extensions]
        for video in video_files:
            fileData['name'] = video.name
            fileData['path'] = str(video.resolve())
            fileData['time'] = video.stat().st_mtime
            fileData['size'] = video.stat().st_size
            fileData['created'] = video.stat().st_ctime
            fileList.append(fileData)
    else:
        return jsonify({"error": f"The directory {'uploads'} does not exist or is not a directory."}), 404
    
    return jsonify(fileList)

# 저장할 경로 설정
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

@app.route('/upload_video', methods=['POST'])
def upload_video():
    # 청크를 받을 임시 파일 생성
    file_name = request.headers.get('File-Name')
    file_name = urllib.parse.unquote(file_name)
    file_path = os.path.join(UPLOAD_FOLDER, file_name)

    # 청크 데이터 가져오기
    with open(file_path, 'ab') as f:  # 'ab' 모드로 파일을 열어 데이터를 추가
        f.write(request.data)
    
    # 현재 받은 청크 수와 전체 청크 수 확인
    chunk_number = int(request.headers.get('Chunk-Number'))
    total_chunks = int(request.headers.get('Total-Chunks'))

    # 모든 청크가 전송된 경우에만 경로를 반환
    if chunk_number == total_chunks:
        return jsonify({"file_path": file_path}), 200
    else:
        return jsonify({"message": f"Chunk {chunk_number} received"}), 200

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
