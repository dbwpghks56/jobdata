import os
import pickle
import re
import shutil
import tkinter as tk
import uuid
from tkinter import filedialog, messagebox

import comtypes.client
import cv2
import easyocr
import numpy as np
from PIL import Image, ImageTk
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# EasyOCR 설정
reader = easyocr.Reader(['en', 'ko'], gpu=True)  # 영어와 한국어 설정

def pptx_to_images(pptx_path, output_dir):
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(pptx_path)
    total_slides = len(presentation.slides)
    
    for i, slide in enumerate(presentation.Slides):
        slide_num = str(i + 1).zfill(len(str(total_slides)))
        slide_name = f"slide_{slide_num}.png"
        slide_path = os.path.join(output_dir, slide_name)
        slide.Export(slide_path, "PNG")
    
    presentation.Close()
    powerpoint.Quit()
    
def extract_text_from_image(image):
    result = reader.readtext(np.array(image))
    text = ' '.join([item[1] for item in result])
    return text

def extract_text_from_ppt(directory, ocr_area=(0, 90, 1080, 600)):
    """
    Extract text from PPT slide images within the specified OCR area.
    """
    slide_texts = []
    
    image_files = sorted(
        [f for f in os.listdir(directory) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff'))]
    )

    for image_file in image_files:
        image_path = os.path.join(directory, image_file)
        image = Image.open(image_path)
        cropped_image = image.crop(ocr_area)  # Use the OCR area specified
        cropped_image.save(image_path)  # Save the cropped image
        text = extract_text_from_image(cropped_image)
        text = re.sub(r'\#\d+', '', text)
        
        slide_texts.append(text)
        
    # Save slide texts to file
    # with open('slide_texts.txt', 'w') as f:
    #     for text in slide_texts:
    #         f.write(text + '\n')
    
    return slide_texts

def process_frame(frame, scale_factor=0.3):
    resized_frame = cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor)
    cropped_frame = resized_frame[60:600, 0:1080]
    gray = cv2.cvtColor(cropped_frame, cv2.COLOR_BGR2GRAY)
    result = reader.readtext(gray)
    text = ' '.join([item[1] for item in result])
    return text

def cache_frame_texts(video_path, frame_skip, scale_factor=0.3, cache_file='frame_texts_cache.pkl'):
    cap = cv2.VideoCapture(video_path)
    frame_texts = []
    frame_images = []
    frame_indices = []
    frame_count = 0
    
    while True:
        success, frame = cap.read()
        if not success:
            break
        if frame_count % frame_skip == 0:
            text = process_frame(frame, scale_factor)
            resized_frame = cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor)
            frame_texts.append(text)
            frame_images.append(resized_frame)
            frame_indices.append(frame_count)
        frame_count += 1
    
    # Save frame texts to file
    # with open('frame_texts.txt', 'w') as f:
    #     for text in frame_texts:
    #         f.write(text + '\n')
    
    with open(cache_file, 'wb') as f:
        pickle.dump((frame_texts, frame_images, frame_indices), f)

def load_cached_frame_texts(cache_file='frame_texts_cache.pkl'):
    with open(cache_file, 'rb') as f:
        return pickle.load(f)
    
def process_video_and_annotate_ppt(ppt_path, video_path, image_directory, output_path, ocr_area, frame_skip=30, cache_file='frame_texts_cache.pkl'):
    slide_texts = extract_text_from_ppt(image_directory, ocr_area=ocr_area)
    
    if not os.path.exists(cache_file):
        cache_frame_texts(video_path, frame_skip, cache_file=cache_file)
    
    frame_texts, frame_images, frame_indices = load_cached_frame_texts(cache_file)
    
    vectorizer = TfidfVectorizer().fit_transform(slide_texts + frame_texts)
    vectors = vectorizer.toarray()
    slide_vectors = vectors[:len(slide_texts)]
    frame_vectors = vectors[len(slide_texts):]
    similarities = cosine_similarity(slide_vectors, frame_vectors)
    
    most_similar_frames = similarities.argmax(axis=1)
    captured_images = [frame_images[i] for i in most_similar_frames]
    
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        max_sim_val = similarities[i][most_similar_frames[i]]
        img_path = image_directory + f"//most_similar_frame_{i + 1}.png"
        cv2.imwrite(img_path, captured_images[i])
        left = Inches(0.5)
        top = Inches(5)
        pic = slide.shapes.add_picture(img_path, Inches(-9.6), Inches(0.3), height=Inches(5), width=Inches(9.5))
        
        text_box = slide.shapes.add_textbox(Inches(11.3), top - Inches(0.5), Inches(2), Inches(0.5))
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.text = f"유사도 : {max_sim_val:.2f}"
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 0)
        
        if max_sim_val <= 0.25:
            rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.3), Inches(2.3), Inches(2), Inches(2.1))
            rectangle.fill.solid()
            rectangle.fill.fore_color.rgb = RGBColor(255, 0, 0)
        
    prs.save(output_path)

def select_ppt_file():
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    ppt_entry.delete(0, tk.END)
    ppt_entry.insert(0, file_path)

def select_video_file():
    file_path = filedialog.askopenfilename(filetypes=[("Video files", "*.mp4;*.avi;*.mov")])
    video_entry.delete(0, tk.END)
    video_entry.insert(0, file_path)

def run_processing():
    ppt_path = ppt_entry.get()
    unique = str(uuid.uuid4()).split('-')[0]
    ppt_folder = os.path.dirname(ppt_path)
    image_directory = ppt_folder + "//slides_" + unique
    video_path = video_entry.get()
    output_path = ppt_folder + "//검수_" + os.path.basename(ppt_path)
    
    if not ppt_path or not image_directory or not video_path or not output_path:
        messagebox.showerror("Error", "모든 파일과 디렉토리를 선택해주세요.")
        return
    
    try:
        # Retrieve the crop values from the GUI
        left = int(crop_left_entry.get())
        top = int(crop_top_entry.get())
        right = int(crop_right_entry.get())
        bottom = int(crop_bottom_entry.get())
        
        pptx_to_images(ppt_path, image_directory)
        process_video_and_annotate_ppt(ppt_path, video_path, image_directory, output_path, ocr_area=(left, top, right, bottom), cache_file=f'frame_texts_cache_{unique}.pkl')
        
        os.remove('frame_texts_cache.pkl')
        shutil.rmtree(image_directory)
        
        messagebox.showinfo("Success", "처리가 완료되었습니다.")
    except Exception as e:
        messagebox.showerror("Error", f"오류 발생: {e}")

def select_image_file():
    file_path = filedialog.askopenfilename(filetypes=[("Image files", "*.png;*.jpg;*.jpeg;*.bmp;*.tiff")])
    image_entry.delete(0, tk.END)
    image_entry.insert(0, file_path)
    
def preview_crop():
    try:
        image_path = image_entry.get()
        if not image_path:
            messagebox.showerror("Error", "이미지 파일을 선택해주세요.")
            return

        image = Image.open(image_path)

        # Crop using the values from the entry boxes
        left = int(crop_left_entry.get())
        top = int(crop_top_entry.get())
        right = int(crop_right_entry.get())
        bottom = int(crop_bottom_entry.get())

        cropped_image = image.crop((left, top, right, bottom))
        
        # Resize the image to fit into the preview area if necessary
        cropped_image.thumbnail((preview_width, preview_height), Image.LANCZOS)
        
        # Display the cropped image
        cropped_image_tk = ImageTk.PhotoImage(cropped_image)
        preview_canvas.delete("all")  # Clear the canvas
        preview_canvas.create_image(0, 0, image=cropped_image_tk, anchor="nw")
        preview_canvas.config(scrollregion=preview_canvas.bbox("all"))  # Update scrollregion
        preview_canvas.image = cropped_image_tk

    except Exception as e:
        messagebox.showerror("Error", f"오류 발생: {e}")

# GUI 설정
root = tk.Tk()
root.title("영상 검수 ( PPT )")

# Define the preview area size
preview_width = 400  # Set desired width for preview area
preview_height = 300  # Set desired height for preview area

tk.Label(root, text="스토리보드(PPT) :").grid(row=0, column=0, padx=5, pady=5)
ppt_entry = tk.Entry(root, width=50)
ppt_entry.grid(row=0, column=1, padx=5, pady=5)
tk.Button(root, text="선택", command=select_ppt_file).grid(row=0, column=2, padx=5, pady=5)

tk.Label(root, text="영상 :", anchor='e').grid(row=1, column=0, padx=5, pady=5)
video_entry = tk.Entry(root, width=50)
video_entry.grid(row=1, column=1, padx=5, pady=5)
tk.Button(root, text="선택", command=select_video_file).grid(row=1, column=2, padx=5, pady=5)

tk.Label(root, text="이미지 파일 :").grid(row=2, column=0, padx=5, pady=5)
image_entry = tk.Entry(root, width=50)
image_entry.grid(row=2, column=1, padx=5, pady=5)
tk.Button(root, text="선택", command=select_image_file).grid(row=2, column=2, padx=5, pady=5)

# Crop dimensions entries
tk.Label(root, text="크롭 영역 (Left, Top, Right, Bottom) :").grid(row=3, column=0, padx=5, pady=5)
crop_left_entry = tk.Entry(root, width=10)
crop_left_entry.grid(row=3, column=1, padx=5, pady=5, sticky="w")
crop_left_entry.insert(0, "0")  # Default value

crop_top_entry = tk.Entry(root, width=10)
crop_top_entry.grid(row=3, column=1, padx=5, pady=5)
crop_top_entry.insert(0, "90")  # Default value

crop_right_entry = tk.Entry(root, width=10)
crop_right_entry.grid(row=3, column=1, padx=5, pady=5, sticky="e")
crop_right_entry.insert(0, "1080")  # Default value

crop_bottom_entry = tk.Entry(root, width=10)
crop_bottom_entry.grid(row=3, column=2, padx=5, pady=5, sticky="w")
crop_bottom_entry.insert(0, "600")  # Default value

tk.Button(root, text="크롭 미리보기", command=preview_crop).grid(row=4, column=0, columnspan=3, pady=10)

# Preview section with Canvas and Scrollbars
preview_frame = tk.Frame(root)
preview_frame.grid(row=5, column=0, columnspan=3, pady=10)

preview_canvas = tk.Canvas(preview_frame, width=preview_width, height=preview_height, bg="white")
preview_canvas.pack(side="left", fill="both", expand=True)

# Add scrollbars
h_scrollbar = tk.Scrollbar(preview_frame, orient="horizontal", command=preview_canvas.xview)
h_scrollbar.pack(side="bottom", fill="x")

v_scrollbar = tk.Scrollbar(preview_frame, orient="vertical", command=preview_canvas.yview)
v_scrollbar.pack(side="right", fill="y")

preview_canvas.config(xscrollcommand=h_scrollbar.set, yscrollcommand=v_scrollbar.set)

tk.Button(root, text="실행", command=run_processing).grid(row=6, column=0, columnspan=3, pady=10)

root.mainloop()

