import os
import pickle
import re
import shutil

import comtypes.client
import cv2
import easyocr
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# EasyOCR 설정
reader = easyocr.Reader(['en', 'ko'], gpu=True)  # 영어와 한국어 설정 및 GPU 사용

def pptx_to_images(pptx_path, output_dir):
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    
    presentation = powerpoint.Presentations.Open(pptx_path)
    total_slides = len(presentation.slides)
    
    for i, slide in enumerate(presentation.Slides):
        slide_num = str(i + 1).zfill(len(str(total_slides)))
        slide_name = f"slide_{slide_num}.png"
        slide_path = os.path.join(output_dir, slide_name)
        slide.Export(slide_path, "PNG")
    
    presentation.Close()
    powerpoint.Quit()

def preprocess_image(image):
    # 이미지를 numpy 배열로 변환
    if not isinstance(image, np.ndarray):
        image = np.array(image)
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    denoised = cv2.fastNlMeansDenoising(gray, None, 30, 7, 21)
    enhanced = cv2.equalizeHist(denoised)
    return enhanced

def extract_text_from_image(image):
    preprocessed = preprocess_image(image)
    result = reader.readtext(preprocessed)
    text = ' '.join([item[1] for item in result])
    return postprocess_text(text)

def postprocess_text(text):
    text = re.sub(r'[^a-zA-Z0-9가-힣\s]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def extract_text_from_ppt(directory, ocr_area=(0, 90, 1080, 600)):
    slide_texts = []
    image_files = sorted([f for f in os.listdir(directory) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff'))])
    
    for image_file in image_files:
        image_path = os.path.join(directory, image_file)
        image = Image.open(image_path)
        cropped_image = image.crop(ocr_area)
        slide_texts.append(extract_text_from_image(cropped_image))
    
    return slide_texts

def process_frame(frame, scale_factor=0.5):
    if frame is None:
        return ""
    resized_frame = cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor)
    preprocessed = preprocess_image(resized_frame)
    result = reader.readtext(preprocessed)
    text = ' '.join([item[1] for item in result])
    return postprocess_text(text)

def cache_frame_texts(video_path, frame_skip, scale_factor=0.5, cache_file='frame_texts_cache.pkl'):
    cap = cv2.VideoCapture(video_path)
    frame_texts = []
    frame_images = []
    frame_indices = []
    frame_count = 30
    
    while True:
        success, frame = cap.read()
        if not success:
            break
        if frame_count % frame_skip == 0:
            text = process_frame(frame, scale_factor)
            if text:  # Only add if text is not empty
                frame_texts.append(text)
                frame_images.append(cv2.resize(frame, (0, 0), fx=scale_factor, fy=scale_factor))
                frame_indices.append(frame_count)
        frame_count += 1
    
    with open(cache_file, 'wb') as f:
        pickle.dump((frame_texts, frame_images, frame_indices), f)

def load_cached_frame_texts(cache_file='frame_texts_cache.pkl'):
    with open(cache_file, 'rb') as f:
        return pickle.load(f)

def process_video_and_annotate_ppt(ppt_path, video_path, image_directory, frame_skip=60, cache_file='frame_texts_cache.pkl'):
    slide_texts = extract_text_from_ppt(image_directory)
    
    if not os.path.exists(cache_file):
        cache_frame_texts(video_path, frame_skip, cache_file=cache_file)
    
    frame_texts, frame_images, frame_indices = load_cached_frame_texts(cache_file)
    
    vectorizer = TfidfVectorizer().fit_transform(slide_texts + frame_texts)
    vectors = vectorizer.toarray()
    slide_vectors = vectors[:len(slide_texts)]
    frame_vectors = vectors[len(slide_texts):]
    similarities = cosine_similarity(slide_vectors, frame_vectors)
    
    most_similar_frames = similarities.argmax(axis=1)
    captured_images = [frame_images[i] for i in most_similar_frames]
    
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        max_sim_val = similarities[i][most_similar_frames[i]]
        img_path = os.path.join(image_directory, f"most_similar_frame_{i + 1}.png")
        cv2.imwrite(img_path, captured_images[i])
        left = Inches(0.5)
        top = Inches(5)
        pic = slide.shapes.add_picture(img_path, left, top, height=Inches(3))
        
        text_box = slide.shapes.add_textbox(left, top - Inches(0.5), Inches(3), Inches(0.5))
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.text = f"Similarity: {max_sim_val:.2f}"
        
        if max_sim_val <= 0.4:
            slide_background = slide.background
            fill = slide_background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)
        
    prs.save("annotated_presentation.pptx")

# 실행 예제
ppt_path = "C://Users//wpghk//Desktop//automatic_text//test.pptx"
image_directory = "C://Users//wpghk//Desktop//automatic_text//slides"
video_path = "C://Users//wpghk//Desktop//automatic_text//(종편본)관광창업론_01주차_02_학습내용_수정07.mp4"

pptx_to_images(ppt_path, image_directory)
process_video_and_annotate_ppt(ppt_path, video_path, image_directory)

# 캐시 파일 삭제
os.remove('frame_texts_cache.pkl')

# 캐시 파일이 삭제되었는지 확인
os.path.exists('frame_texts_cache.pkl')

# image_directory 삭제
shutil.rmtree(image_directory)
